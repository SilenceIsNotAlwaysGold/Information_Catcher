# -*- coding: utf-8 -*-
"""AI PPT —— 输入主题 → AI 生大纲 JSON → python-pptx 渲染 .pptx。

计费：
  ppt_outline   — 生成大纲（一次 LLM 调用）
  （渲染本身不耗 AI；不计费）

后续可加：上传用户已有 .pptx 让 AI 改造（按用户指令重写某页/扩充某页/换风格）。
"""
from __future__ import annotations

import io
import json
import logging
import os
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

from ..auth import get_current_user, get_current_user_flex
from ...services import ai_client, monitor_db, storage
from ...services import db as _db

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/studio/ppt", tags=["AI-Studio-PPT"])
DB_PATH = monitor_db.DB_PATH

# 渲染产物目录：必须持久（系统临时目录会被 OS/重启清理 → 列表显示"已渲染"
# 但下载 404）。落到项目 database/ 下，与 ppt_templates 同基准。
_PPT_CACHE_DIR = os.path.join(os.path.dirname(str(DB_PATH)), "ppt_cache")
os.makedirs(_PPT_CACHE_DIR, exist_ok=True)


# ── DB helpers ───────────────────────────────────────────────────────────

async def _get_project(pid: int, user_id: int):
    async with _db.connect(DB_PATH) as db:
        db.row_factory = _db.Row
        async with db.execute(
            "SELECT * FROM ppt_projects WHERE id=? AND user_id=?", (int(pid), int(user_id)),
        ) as cur:
            row = await cur.fetchone()
    return dict(row) if row else None


async def _touch(pid: int, **fields):
    sets = ["updated_at=datetime('now','localtime')"]
    vals: List[Any] = []
    for k, v in fields.items():
        sets.append(f"{k}=?"); vals.append(v)
    vals.append(int(pid))
    async with _db.connect(DB_PATH) as db:
        await db.execute(f"UPDATE ppt_projects SET {', '.join(sets)} WHERE id=?", vals)
        await db.commit()


# ── AI 生大纲 ────────────────────────────────────────────────────────────

_OUTLINE_SYSTEM = """你是 PPT 内容策划。根据用户给的主题、目标页数、观众和风格，输出大纲。

严格只输出 JSON（不要任何代码块标记或解释），格式：
{
  "title":"整套 PPT 的总标题",
  "pages":[
    {"layout":"cover|bullets|two_column|image_right|quote|closing",
     "title":"页标题",
     "bullets":["要点1","要点2","要点3"],
     "left":["左栏要点（仅 two_column 用）"],
     "right":["右栏要点（仅 two_column 用）"],
     "image_query":"配图英文关键词（image_right / cover 用，比如 'team meeting strategy'）",
     "quote":"金句正文（仅 quote 布局用）",
     "quote_author":"金句作者（quote 布局用，可空）"},
    ...
  ]
}

layout 用法：
- cover    — 第 1 页必须用，整页大标题，可配 image_query 作背景意象
- bullets  — 标准要点页，3-5 条 bullets
- two_column — 对比 / 两面观点 / 之前 vs 之后；填 left/right 数组各 3-5 条
- image_right — 左文右图，bullets 3 条以内，必须给 image_query
- quote    — 整页一句金句 / 数据冲击，填 quote + quote_author（不要 bullets）
- closing  — 结尾页，标题如「谢谢」「Q&A」「联系方式」，bullets 0-3 条

要点：
- 第 1 页用 cover，最后一页用 closing
- 总页数 = 用户给的 target_pages（包含 cover 和 closing）
- bullets 每条 12-30 字，言之有物；不要"概述/介绍/总结"这种空话
- image_query 用**英文** 2-5 个词，方便后续 Pexels 搜图
"""


class CreateProjectIn(BaseModel):
    title: str = ""
    topic: str = Field(..., min_length=2, max_length=500)
    target_pages: int = Field(10, ge=3, le=30)
    style_hint: str = ""
    audience: str = ""
    text_model_id: Optional[int] = None
    template_id: Optional[int] = None       # v2.1: 用户上传的 .pptx 风格模板
    image_source: str = "none"              # v2.1: none | pexels | ai


def _strip_codeblock(s: str) -> str:
    t = s.strip()
    if t.startswith("```"):
        t = t.strip("`")
        if t.lower().startswith("json"):
            t = t[4:].strip()
    return t


@router.post("/projects", summary="新建项目 + AI 生大纲（扣 ppt_outline 点）")
async def create_project(body: CreateProjectIn, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    user_prompt = (
        f"PPT 主题：{body.topic.strip()}\n"
        f"目标页数：{body.target_pages} 页\n"
        f"目标观众：{body.audience or '通用'}\n"
        f"风格：{body.style_hint or '商务简洁'}\n\n"
        f"请输出大纲 JSON。"
    )
    raw = await ai_client.call_text(
        user_prompt, model_id=body.text_model_id, user_id=uid,
        feature="ppt_outline", system_prompt=_OUTLINE_SYSTEM,
        temperature=0.7, max_tokens=3000,
        task_ref=ai_client.make_task_ref(
            "ppt_outline", uid, body.topic, body.target_pages,
            body.audience or "", body.style_hint or "",
        ),
        expect_json=True,
    )
    txt = _strip_codeblock(raw)
    try:
        plan = json.loads(txt)
        if not isinstance(plan, dict) or not isinstance(plan.get("pages"), list) or not plan["pages"]:
            raise ValueError("缺 pages")
    except Exception as e:
        raise HTTPException(
            502,
            f"AI 返回的大纲不是合法 JSON：{str(e)[:100]}；原始（前 200 字符）：{raw[:200]}",
        )
    title = (body.title or plan.get("title") or body.topic).strip()[:100]
    img_src = body.image_source if body.image_source in ("none", "pexels", "ai") else "none"
    async with _db.connect(DB_PATH) as db:
        cur = await db.execute(
            "INSERT INTO ppt_projects(user_id, title, topic, target_pages, style_hint, audience, "
            "plan_json, status, text_model_id, template_id, image_source) "
            "VALUES (?,?,?,?,?,?,?, 'outlined', ?,?,?)",
            (uid, title, body.topic.strip()[:500], int(body.target_pages),
             body.style_hint.strip()[:80], body.audience.strip()[:80],
             json.dumps(plan, ensure_ascii=False), body.text_model_id,
             body.template_id, img_src),
        )
        await db.commit()
        pid = cur.lastrowid
    return {"ok": True, "id": pid, "title": title, "plan": plan}


def _parse_pptx_to_plan(blob: bytes) -> Dict[str, Any]:
    """解析用户上传的 .pptx → 我们的 plan JSON {title, pages:[{title,bullets[]}]}。

    规则（简单粗暴够用）：
      - 第一张幻灯片 → 取最长的非空文本块当 plan.title
      - 后续每张 → 取第一个非空文本块当 page.title，剩下所有文本块切行 → bullets
      - 空白幻灯片跳过
    """
    from pptx import Presentation
    prs = Presentation(io.BytesIO(blob))
    pages: List[Dict[str, Any]] = []
    cover_title = ""

    import re as _re
    _BULLET_RE = _re.compile(r"^[•●◦▪·\-\*•·▪◦●○]\s*")
    _PAGENUM_RE = _re.compile(r"^\s*\d{1,3}\s*/\s*\d{1,3}\s*$")  # "01 / 30"

    def _slide_texts(slide) -> List[str]:
        """收集 slide 内所有文字。按 shape 位置 (top, left) 排序，使阅读顺序近似从上到下。
        同一 shape 内多段落各算一行；去前缀 bullet 符号；丢弃 "01/30" 这种纯页码。"""
        items = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            top = shape.top if shape.top is not None else 0
            left = shape.left if shape.left is not None else 0
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text or "" for run in para.runs).strip()
                if not line:
                    line = (para.text or "").strip()
                if not line:
                    continue
                line = _BULLET_RE.sub("", line).strip()
                if not line or _PAGENUM_RE.match(line):
                    continue
                items.append((top, left, line))
        items.sort(key=lambda x: (x[0], x[1]))
        return [t[2] for t in items]

    for idx, slide in enumerate(prs.slides):
        texts = _slide_texts(slide)
        if not texts:
            continue
        if idx == 0 and not pages:
            # 把最长的当封面 title
            cover_title = max(texts, key=len)[:100]
            # 如果还有别的文本块，剩下的当第 1 页的 bullets（避免完全丢内容）
            others = [t for t in texts if t != cover_title]
            if others:
                pages.append({"title": cover_title or f"Slide 1", "bullets": others[:8]})
            continue
        page_title = texts[0][:80]
        bullets = [t for t in texts[1:] if t and t != page_title][:8]
        pages.append({"title": page_title, "bullets": bullets})

    if not pages:
        raise ValueError("PPT 里没解析到任何文字内容")
    return {"title": cover_title or pages[0]["title"], "pages": pages}


@router.post("/projects/upload", summary="上传 .pptx 创建项目（解析成大纲）")
async def upload_pptx(
    file: UploadFile = File(...),
    title: str = Form(""),
    style_hint: str = Form(""),
    current_user: dict = Depends(get_current_user),
):
    """不扣点 —— 只是解析现有 .pptx 文件不动 AI。AI 改造在 /revise 里扣。"""
    uid = int(current_user["id"])
    name = (file.filename or "").lower()
    if not name.endswith(".pptx"):
        raise HTTPException(400, "只支持 .pptx 格式（不支持老 .ppt / WPS .dps）")
    blob = await file.read()
    if len(blob) > 30 * 1024 * 1024:
        raise HTTPException(400, "文件超过 30 MB")
    try:
        plan = _parse_pptx_to_plan(blob)
    except Exception as e:
        raise HTTPException(400, f"解析 .pptx 失败：{type(e).__name__}: {str(e)[:160]}")
    plan_title = (title.strip() or plan.get("title") or "导入的 PPT")[:100]
    plan["title"] = plan_title
    pages_n = len(plan["pages"])
    topic = f"(导入自 {file.filename or 'uploaded.pptx'})"
    async with _db.connect(DB_PATH) as db:
        cur = await db.execute(
            "INSERT INTO ppt_projects(user_id, title, topic, target_pages, style_hint, audience, "
            "plan_json, status) VALUES (?,?,?,?,?,?,?, 'outlined')",
            (uid, plan_title, topic[:500], pages_n,
             style_hint.strip()[:80], "",
             json.dumps(plan, ensure_ascii=False)),
        )
        await db.commit()
        pid = cur.lastrowid
    return {"ok": True, "id": pid, "title": plan_title, "plan": plan, "pages": pages_n}


_REVISE_SYSTEM = """你是 PPT 内容修订助手。
用户给你：现有的 PPT 大纲 JSON + 一段修改指令。
你需要返回**修改后**的完整大纲 JSON（保留未涉及的页面原样）。
严格只输出 JSON（不加代码块标记或解释），结构同生大纲：
  {"title":"...","pages":[{"layout":"...","title":"...","bullets":[...],"left":[...],"right":[...],"image_query":"...","quote":"...","quote_author":"..."}, ...]}
layout 枚举：cover | bullets | two_column | image_right | quote | closing
要点：
- 用户没指定的页面，原样保留（含 layout 字段）；
- 用户说"加一页 XX"则在合理位置插入新页，自己选合适 layout；
- 用户说"删掉第N页"则真的删除；
- bullets 每条 12-30 字，言之有物。
"""


class ReviseIn(BaseModel):
    instruction: str = Field(..., min_length=2, max_length=500)
    text_model_id: Optional[int] = None


@router.post("/projects/{pid}/revise", summary="按指令让 AI 修改大纲（扣 ppt_outline 点）")
async def revise_project(
    pid: int, body: ReviseIn,
    current_user: dict = Depends(get_current_user),
):
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    try:
        plan = json.loads(p.get("plan_json") or "{}")
    except Exception:
        raise HTTPException(400, "现有大纲 JSON 异常")
    if not plan.get("pages"):
        raise HTTPException(400, "现有大纲为空，无法修订；请先生成或上传")

    user_prompt = (
        f"=== 现有大纲（JSON）===\n{json.dumps(plan, ensure_ascii=False)}\n\n"
        f"=== 修改指令 ===\n{body.instruction.strip()}\n\n"
        f"请输出修改后的完整大纲 JSON。"
    )
    raw = await ai_client.call_text(
        user_prompt,
        model_id=body.text_model_id or p.get("text_model_id"),
        user_id=uid, feature="ppt_outline",
        system_prompt=_REVISE_SYSTEM,
        temperature=0.4, max_tokens=4000,
        task_ref=ai_client.make_task_ref("ppt_revise", pid, body.instruction),
        expect_json=True,
    )
    txt = _strip_codeblock(raw)
    try:
        new_plan = json.loads(txt)
        if not isinstance(new_plan, dict) or not isinstance(new_plan.get("pages"), list) or not new_plan["pages"]:
            raise ValueError("缺 pages")
    except Exception as e:
        raise HTTPException(
            502,
            f"AI 返回的修订大纲不是合法 JSON：{str(e)[:100]}；原始（前 200 字符）：{raw[:200]}",
        )
    # 修订完旧 .pptx 文件作废
    fpath = os.path.join(_PPT_CACHE_DIR, f"ppt_{uid}_{pid}.pptx")
    try:
        if os.path.exists(fpath):
            os.remove(fpath)
    except Exception:
        pass
    await _touch(
        pid,
        plan_json=json.dumps(new_plan, ensure_ascii=False),
        target_pages=len(new_plan["pages"]),
        status="outlined",
        pptx_url="",
    )
    return {"ok": True, "id": pid, "plan": new_plan, "pages": len(new_plan["pages"])}


@router.get("/projects", summary="我的 PPT 项目列表")
async def list_projects(current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    async with _db.connect(DB_PATH) as db:
        db.row_factory = _db.Row
        async with db.execute(
            "SELECT id, title, topic, target_pages, style_hint, audience, status, "
            "pptx_url, template_id, COALESCE(image_source,'none') AS image_source, "
            "created_at, updated_at "
            "FROM ppt_projects WHERE user_id=? ORDER BY updated_at DESC, id DESC",
            (uid,),
        ) as cur:
            rows = [dict(r) for r in await cur.fetchall()]
    return {"projects": rows}


@router.get("/projects/{pid}", summary="项目详情（含大纲 JSON）")
async def get_project_detail(pid: int, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    try:
        p["plan"] = json.loads(p.get("plan_json") or "{}")
    except Exception:
        p["plan"] = {}
    p.pop("plan_json", None)
    return p


class UpdatePlanIn(BaseModel):
    plan: Dict[str, Any]


@router.put("/projects/{pid}/plan", summary="手工编辑大纲 JSON")
async def update_plan(pid: int, body: UpdatePlanIn, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    if not isinstance(body.plan, dict) or not isinstance(body.plan.get("pages"), list):
        raise HTTPException(400, "plan 必须是 {title, pages:[{title,bullets[]}]}")
    await _touch(pid, plan_json=json.dumps(body.plan, ensure_ascii=False))
    return {"ok": True}


@router.delete("/projects/{pid}", summary="删除 PPT 项目")
async def delete_project(pid: int, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    async with _db.connect(DB_PATH) as db:
        await db.execute("DELETE FROM ppt_projects WHERE id=?", (int(pid),))
        await db.commit()
    return {"ok": True}


# ── 渲染 .pptx ───────────────────────────────────────────────────────────

def _build_pptx(
    plan: Dict[str, Any],
    style_hint: str = "",
    *,
    template_path: Optional[str] = None,
    images_by_idx: Optional[Dict[int, bytes]] = None,
) -> bytes:
    """用 python-pptx 把大纲渲染成 .pptx 二进制。

    template_path 非空 → 用它作 master/theme/layouts 的底版（保留公司 VI 风格），
                          先清空模板带的样例幻灯片，再用模板的 layouts 重塞内容。
    template_path 为空 → 沿用我们内置的双色配色（商务深蓝 / 极简白）。
    images_by_idx     → {page_idx: jpg/png bytes}；image_right / cover 布局优先用
                          这里的图，没有则降级到占位卡片。
                          page_idx 0 = 封面，1+ = pages 数组里的 i+1。
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    if template_path:
        try:
            return _build_from_template(plan, template_path, images_by_idx=images_by_idx)
        except Exception as e:
            # 模板渲染走 python-pptx 私有 API 删 sample slides，不同来源 .pptx
            # （WPS 导出 / 含 sectionLst / notesSlide rels）易触发异常并产出
            # partname 重复的损坏 .pptx（PowerPoint 打开报修复）。任何失败一律
            # 回退内置样式，保证用户拿到的是合法可打开的 .pptx 而非坏文件。
            logger.warning(
                "[ppt] 模板渲染失败，回退内置样式（避免产损坏 .pptx）: %s", e
            )
            return _build_pptx(
                plan, style_hint, template_path=None, images_by_idx=images_by_idx,
            )

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 配色（双色基底：商务深蓝 / 极简白底）— 这套 palette 传给每个 layout 函数
    is_dark = ("商务" in style_hint) or ("严肃" in style_hint) or ("dark" in style_hint.lower())
    palette = {
        "bg": RGBColor(0x10, 0x21, 0x3F) if is_dark else RGBColor(0xFF, 0xFF, 0xFF),
        "fg": RGBColor(0xFF, 0xFF, 0xFF) if is_dark else RGBColor(0x18, 0x20, 0x33),
        "accent": RGBColor(0x4F, 0x9C, 0xF5),
        "subtle": RGBColor(0xAA, 0xBB, 0xD8) if is_dark else RGBColor(0x55, 0x66, 0x77),
    }

    pages = plan.get("pages") or []
    total = len(pages)
    images = images_by_idx or {}

    # 自动加封面：若 plan.pages 没显式 cover，就用 plan.title 注入一张
    if not pages or (pages[0].get("layout") not in ("cover",)):
        _render_layout(prs, palette, {
            "layout": "cover",
            "title": plan.get("title") or "PPT",
            "bullets": [style_hint] if style_hint else [],
        }, page_idx=0, total=total + (0 if pages and pages[0].get("layout") == "cover" else 1),
        image_bytes=images.get(0))

    # 渲染所有 page；page_idx 从 1 开始（封面占 0）
    for i, pg in enumerate(pages, start=1):
        _render_layout(prs, palette, pg, page_idx=i, total=total + 1,
                       image_bytes=images.get(i))

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


# ── 单页布局渲染器 ─────────────────────────────────────────────────────

def _render_layout(prs, palette, page: Dict[str, Any], *, page_idx: int, total: int,
                   image_bytes: Optional[bytes] = None):
    """按 page.layout 派发到对应布局函数；未识别的 layout 回退到 bullets。"""
    layout = (page.get("layout") or "bullets").strip().lower()
    funcs = {
        "cover": _layout_cover,
        "bullets": _layout_bullets,
        "two_column": _layout_two_column,
        "image_right": _layout_image_right,
        "quote": _layout_quote,
        "closing": _layout_closing,
    }
    fn = funcs.get(layout, _layout_bullets)
    fn(prs, palette, page, page_idx=page_idx, total=total, image_bytes=image_bytes)


def _blank_slide(prs, palette):
    """加一张空白底，铺纯色背景。返回 slide 对象。所有布局都从这里起步。"""
    from pptx.util import Inches
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = palette["bg"]
    return slide


def _add_page_number(slide, palette, page_idx: int, total: int):
    """页码（左下角，accent 色，小字）— 封面和 quote 用大字时不加。"""
    from pptx.util import Inches, Pt
    if page_idx <= 0:
        return
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(2), Inches(0.4))
    tf = tx.text_frame
    tf.text = f"{page_idx:02d} / {total:02d}"
    tf.paragraphs[0].runs[0].font.size = Pt(12)
    tf.paragraphs[0].runs[0].font.color.rgb = palette["accent"]


def _add_title(slide, palette, text: str, *, top_inches: float = 0.5,
               size: int = 36, with_bar: bool = True):
    """给页加大标题（左上角，可选左侧 accent 装饰条）。"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    if with_bar:
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(top_inches + 0.1), Inches(0.12), Inches(0.8))
        bar.fill.solid()
        bar.fill.fore_color.rgb = palette["accent"]
        bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.7), Inches(top_inches), Inches(12), Inches(1.0))
    tf = tx.text_frame
    tf.text = (text or "").strip()[:120]
    if not tf.paragraphs[0].runs:
        return
    run = tf.paragraphs[0].runs[0]
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = palette["fg"]


def _add_bullets(slide, palette, bullets: List[Any], *,
                 left: float, top: float, width: float, height: float,
                 font_size: int = 22, max_n: int = 8):
    """通用 bullets 文本框。"""
    from pptx.util import Inches, Pt
    items = [str(b).strip() for b in (bullets or []) if str(b).strip()][:max_n]
    if not items:
        return
    tx = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.space_after = Pt(12)
        for r in p.runs:
            r.font.size = Pt(font_size)
            r.font.color.rgb = palette["fg"]


def _layout_cover(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """封面：整页大标题居中 + 可选 subtitle（bullets[0]）；有 image_bytes 时铺满底图 + 暗罩。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide = _blank_slide(prs, palette)
    # 全屏底图（如果有）
    if image_bytes:
        try:
            from io import BytesIO
            from pptx.enum.shapes import MSO_SHAPE
            slide.shapes.add_picture(BytesIO(image_bytes),
                Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
            # 暗罩：黑色半透矩形，让文字醒目。
            # 注意：python-pptx 的 FillFormat 没有可写的 transparency 属性，
            # 直接赋值会 AttributeError 被外层 except 吞掉 → 矩形仍是 100%
            # 不透明纯黑，把底图整张盖死（旧 bug）。改为往 a:srgbClr 注入
            # <a:alpha> 实现真半透（OOXML alpha：100000=完全不透明）。
            from pptx.oxml.ns import qn
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0), Inches(13.333), Inches(7.5))
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            overlay.line.fill.background()
            _srgb = overlay._element.spPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
            _srgb.append(_srgb.makeelement(qn("a:alpha"), {"val": "42000"}))  # 42% 不透明，底图透出 ~58%
        except Exception:
            pass
    title = (page.get("title") or "PPT").strip()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.2))
    tf = tx.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = 1  # center
    if p.runs:
        run = p.runs[0]
        run.font.size = Pt(60)
        run.font.bold = True
        run.font.color.rgb = palette["fg"]
    bullets = page.get("bullets") or []
    subtitle = bullets[0] if bullets else ""
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(0.6))
        sf = sub.text_frame
        sf.text = str(subtitle).strip()
        sf.paragraphs[0].alignment = 1
        if sf.paragraphs[0].runs:
            sf.paragraphs[0].runs[0].font.size = Pt(20)
            sf.paragraphs[0].runs[0].font.color.rgb = palette["accent"]


def _layout_bullets(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """标准要点页（bullets 布局不用图，ignore image_bytes）。"""
    slide = _blank_slide(prs, palette)
    _add_page_number(slide, palette, page_idx, total)
    _add_title(slide, palette, page.get("title") or "")
    _add_bullets(slide, palette, page.get("bullets") or [],
                 left=0.7, top=1.8, width=12, height=5.0)


def _layout_two_column(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """对比页：左右两栏，各 3-5 条 bullets。两栏布局不用图。"""
    from pptx.util import Inches, Pt
    slide = _blank_slide(prs, palette)
    _add_page_number(slide, palette, page_idx, total)
    _add_title(slide, palette, page.get("title") or "")

    # 左栏 / 右栏 — 用 left/right 字段；没有就把 bullets 平均切两半作 fallback
    left_items = page.get("left") or []
    right_items = page.get("right") or []
    if not (left_items or right_items):
        bullets = page.get("bullets") or []
        mid = (len(bullets) + 1) // 2
        left_items, right_items = bullets[:mid], bullets[mid:]

    # 列头（如果有 left_title / right_title 字段）
    def _col_header(text, x):
        if not text:
            return
        tx = slide.shapes.add_textbox(Inches(x), Inches(1.7), Inches(5.8), Inches(0.5))
        tf = tx.text_frame
        tf.text = str(text).strip()
        if tf.paragraphs[0].runs:
            tf.paragraphs[0].runs[0].font.size = Pt(18)
            tf.paragraphs[0].runs[0].font.bold = True
            tf.paragraphs[0].runs[0].font.color.rgb = palette["accent"]

    _col_header(page.get("left_title"), 0.7)
    _col_header(page.get("right_title"), 6.8)

    _add_bullets(slide, palette, left_items,
                 left=0.7, top=2.3, width=5.8, height=4.8, font_size=18, max_n=5)
    _add_bullets(slide, palette, right_items,
                 left=6.8, top=2.3, width=5.8, height=4.8, font_size=18, max_n=5)

    # 中间分隔竖线
    from pptx.enum.shapes import MSO_SHAPE
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(2.0), Inches(0.03), Inches(5.0))
    line.fill.solid()
    line.fill.fore_color.rgb = palette["subtle"]
    line.line.fill.background()


def _layout_image_right(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """左文右图：右侧 4.8 寸预留给图片。
    有 image_bytes → 嵌入；没有 → 占位卡片显示 image_query 提示。"""
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from io import BytesIO
    slide = _blank_slide(prs, palette)
    _add_page_number(slide, palette, page_idx, total)
    _add_title(slide, palette, page.get("title") or "")
    # 左侧 bullets（窄一点，给图腾空间）
    _add_bullets(slide, palette, page.get("bullets") or [],
                 left=0.7, top=1.8, width=6.8, height=5.0, font_size=20, max_n=5)
    # 右侧图片区
    if image_bytes:
        try:
            slide.shapes.add_picture(
                BytesIO(image_bytes),
                Inches(8.0), Inches(1.8), width=Inches(4.8), height=Inches(5.0),
            )
            return
        except Exception:
            pass  # 图坏了 → 走占位
    img_query = (page.get("image_query") or "").strip()
    placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.8), Inches(4.8), Inches(5.0))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = palette["subtle"]
    placeholder.line.fill.background()
    txt = slide.shapes.add_textbox(Inches(8.0), Inches(4.0), Inches(4.8), Inches(0.8))
    tf = txt.text_frame
    tf.text = f"[配图：{img_query or '尚未指定关键词'}]"
    tf.paragraphs[0].alignment = 1
    if tf.paragraphs[0].runs:
        tf.paragraphs[0].runs[0].font.size = Pt(14)
        tf.paragraphs[0].runs[0].font.color.rgb = palette["bg"]
        tf.paragraphs[0].runs[0].font.italic = True


def _layout_quote(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """金句页：整页一句大字 + 作者署名。"""
    from pptx.util import Inches, Pt
    slide = _blank_slide(prs, palette)
    quote = (page.get("quote") or page.get("title") or "").strip()
    author = (page.get("quote_author") or "").strip()
    tx = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.text = f"“{quote}”"
    tf.paragraphs[0].alignment = 1
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(40)
        run.font.italic = True
        run.font.color.rgb = palette["fg"]
    if author:
        a = slide.shapes.add_textbox(Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.6))
        af = a.text_frame
        af.text = f"— {author}"
        af.paragraphs[0].alignment = 1
        if af.paragraphs[0].runs:
            af.paragraphs[0].runs[0].font.size = Pt(18)
            af.paragraphs[0].runs[0].font.color.rgb = palette["accent"]


def _layout_closing(prs, palette, page, *, page_idx: int, total: int, image_bytes: Optional[bytes] = None):
    """结尾页：大标题（"谢谢"/"Q&A"），可选 0-3 条联系方式式 bullets。"""
    from pptx.util import Inches, Pt
    slide = _blank_slide(prs, palette)
    title = (page.get("title") or "谢谢").strip()
    tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.5))
    tf = tx.text_frame
    tf.text = title
    tf.paragraphs[0].alignment = 1
    if tf.paragraphs[0].runs:
        run = tf.paragraphs[0].runs[0]
        run.font.size = Pt(72)
        run.font.bold = True
        run.font.color.rgb = palette["fg"]
    bullets = [str(b).strip() for b in (page.get("bullets") or []) if str(b).strip()][:3]
    if bullets:
        sub = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(11.7), Inches(1.6))
        sf = sub.text_frame
        sf.word_wrap = True
        for i, b in enumerate(bullets):
            p = sf.paragraphs[0] if i == 0 else sf.add_paragraph()
            p.text = b
            p.alignment = 1
            for r in p.runs:
                r.font.size = Pt(18)
                r.font.color.rgb = palette["accent"]


def _build_from_template(plan: Dict[str, Any], template_path: str, *,
                          images_by_idx: Optional[Dict[int, bytes]] = None) -> bytes:
    """基于用户上传的 .pptx 模板渲染。

    思路：
      1) 打开模板 → 拷贝它的 master / theme / slide_layouts 自动随之
      2) 清掉模板自带的 sample slides（保留 layouts，不清 layouts）
      3) 每页选合适 layout：
           - 第 1 页（封面）→ 优先 layouts[0]（Title Slide）
           - 内容页       → 优先 layouts[1]（Title and Content）
         填它的 placeholder（title / body）；body 没占位符就 add_textbox 兜底
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from copy import deepcopy

    prs = Presentation(template_path)

    # 清掉模板里的 sample slides（保留 layouts/master 不动）
    # python-pptx 没现成 delete_slide API：要同时拆 (1) presentation.xml 里 sldIdLst 的引用、
    # (2) presentation.xml.rels 里指向 slide 的 rel、(3) Package 里的 slide part 本身。
    # 缺第 (2)(3) 步会产生 "Duplicate name: ppt/slides/slide1.xml" 警告 + 文件膨胀。
    prs_part = prs.part
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    sld_ids = list(sldIdLst)
    for sld in sld_ids:
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        # 取 rel → 取 slide_part → 从 package 删
        try:
            rel = prs_part.rels[rId]
            slide_part = rel.target_part
            prs_part.drop_rel(rId)
            # 同时从 package 里抹掉这个 part 的 partname，让后续 add_slide 重用 slide1.xml 不冲突
            try:
                pkg = prs_part.package
                if slide_part.partname in pkg._parts:  # noqa: SLF001
                    del pkg._parts[slide_part.partname]  # noqa: SLF001
            except Exception:
                pass
        except Exception:
            pass
        sldIdLst.remove(sld)

    layouts = prs.slide_layouts
    cover_layout = layouts[0] if len(layouts) > 0 else layouts[0]
    body_layout = layouts[1] if len(layouts) > 1 else cover_layout

    def _fill_title(slide, text: str) -> bool:
        """优先填 title placeholder；找不到返回 False。"""
        try:
            if slide.shapes.title is not None:
                slide.shapes.title.text = text[:120]
                return True
        except Exception:
            pass
        return False

    def _fill_body(slide, lines: List[str]) -> bool:
        """优先填 body/content placeholder（一般 idx>=1，type 不是 title）。"""
        for ph in slide.placeholders:
            try:
                if ph.placeholder_format.idx == 0:  # title 跳过
                    continue
                tf = ph.text_frame
                tf.clear()
                for i, line in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(line).strip()[:200]
                    p.level = 0
                return True
            except Exception:
                continue
        return False

    pages = plan.get("pages") or []
    # ─ 封面页 ─
    cover_slide = prs.slides.add_slide(cover_layout)
    cover_text = (plan.get("title") or "PPT").strip()[:120]
    if not _fill_title(cover_slide, cover_text):
        # layout 里没 title 占位 → 自己 add_textbox
        tx = cover_slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(11.7), Inches(2.0))
        tf = tx.text_frame
        tf.text = cover_text
        tf.paragraphs[0].runs[0].font.size = Pt(48)
        tf.paragraphs[0].runs[0].font.bold = True

    # ─ 内容页 ─
    for pg in pages:
        slide = prs.slides.add_slide(body_layout)
        title_text = (pg.get("title") or "").strip()[:120]
        if title_text:
            if not _fill_title(slide, title_text):
                tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1.0))
                tx.text_frame.text = title_text
                tx.text_frame.paragraphs[0].runs[0].font.size = Pt(32)
                tx.text_frame.paragraphs[0].runs[0].font.bold = True
        bullets = [str(b).strip() for b in (pg.get("bullets") or []) if str(b).strip()][:8]
        if bullets:
            if not _fill_body(slide, bullets):
                # 没 body 占位符 → 自己 add_textbox
                tx = slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(12), Inches(5.4))
                tf = tx.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = f"•  {b}"
                    for r in p.runs:
                        r.font.size = Pt(20)
                    p.space_after = Pt(10)

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()


async def _fetch_images_for_plan(plan: Dict[str, Any], *, source: str, user_id: int) -> Dict[int, bytes]:
    """根据 plan.pages[*].image_query 抓图（cover / image_right 布局优先）。

    source='pexels' → Pexels API
    source='ai'     → ai_client.call_image（按用户的图像模型走，会扣 image 点）

    返回 {page_idx: bytes}；page_idx 与 _build_pptx 内部编号一致：
      0 = 自动/手动 cover, 1+ = plan.pages[i-1]
    """
    from ...services import pexels_client
    pages = plan.get("pages") or []
    # 算 cover 是否自动注入：判断 pages[0].layout 是不是 cover
    auto_cover = not (pages and (pages[0].get("layout") or "").lower() == "cover")
    # 收集每个 page_idx 的查询词
    targets: List[tuple] = []  # (page_idx, query)
    if auto_cover:
        # 自动封面用 plan.title 作 query（英文标题更有效；但中文也照用，Pexels 会模糊匹配）
        q = (plan.get("title") or "").strip()
        if q:
            targets.append((0, q))
    for i, pg in enumerate(pages, start=1):
        layout = (pg.get("layout") or "").lower()
        if layout not in ("cover", "image_right"):
            continue
        q = (pg.get("image_query") or pg.get("title") or "").strip()
        if not q:
            continue
        targets.append((i, q))
    if not targets:
        return {}

    out: Dict[int, bytes] = {}
    if source == "pexels":
        for idx, q in targets:
            got = await pexels_client.search_one(q)
            if got:
                out[idx] = got[0]
    elif source == "ai":
        # AI 生图 — 留给 task #17（C）填具体实现
        for idx, q in targets:
            try:
                blob = await _ai_image_for_ppt(q, user_id=user_id)
                if blob:
                    out[idx] = blob
            except Exception as exc:
                logger.warning("[ppt] AI image 失败 idx=%s q=%s: %s", idx, q, exc)
                # 兜底：尝试 Pexels
                got = await pexels_client.search_one(q)
                if got:
                    out[idx] = got[0]
    return out


async def _ai_image_for_ppt(query: str, *, user_id: int) -> Optional[bytes]:
    """AI 生图：把 image_query → 漂亮的横版 PPT 配图。

    走 ai_client.call_image（按 feature='image' 扣点）；返回 PNG/JPG 二进制，
    失败返回 None，由上层 _fetch_images_for_plan 回退到 Pexels。
    """
    import base64
    q = (query or "").strip()
    if not q:
        return None
    # 改写成更适合做横版 PPT 配图的 prompt
    prompt = (
        f"A clean, professional photograph illustrating: {q}. "
        f"Wide horizontal composition, soft natural lighting, modern editorial style, "
        f"high quality, suitable as a slide deck illustration. No text, no logos."
    )
    try:
        b64s = await ai_client.call_image(
            prompt, user_id=user_id, feature="image", n=1, size="1024x768",
            task_ref=ai_client.make_task_ref("ppt_img", user_id, q),
        )
        if not b64s:
            return None
        return base64.b64decode(b64s[0])
    except Exception as exc:
        # InsufficientCredits 单独透出，让前端能提示余额；其它当成本次失败回退到 Pexels
        from ...services.billing_service import InsufficientCredits
        if isinstance(exc, InsufficientCredits):
            raise
        logger.warning("[ppt] AI 生图失败 q=%s: %s", q, exc)
        return None


@router.post("/projects/{pid}/render", summary="把大纲渲染成 .pptx 文件")
async def render_project(pid: int, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    try:
        plan = json.loads(p.get("plan_json") or "{}")
    except Exception:
        raise HTTPException(400, "大纲 JSON 异常")
    if not plan.get("pages"):
        raise HTTPException(400, "大纲为空，先调用 create_project 或 PUT /plan 设大纲")
    await _touch(pid, status="rendering")
    # v2.1: 用户绑了模板则走模板渲染（公司 VI 风格底版）
    template_path: Optional[str] = None
    if p.get("template_id"):
        from . import ppt_templates as _tpl
        template_path = await _tpl.get_template_path(p["template_id"], uid)
        if not template_path:
            logger.warning(f"[ppt] template_id={p.get('template_id')} not found / file missing — fallback 默认风格")

    # v2.1: 配图来源（pexels / ai / none）
    image_source = (p.get("image_source") or "none").strip().lower()
    images_by_idx: Dict[int, bytes] = {}
    if image_source in ("pexels", "ai"):
        images_by_idx = await _fetch_images_for_plan(plan, source=image_source, user_id=uid)
        if images_by_idx:
            logger.info(f"[ppt] pid={pid} 抓图 {len(images_by_idx)} 张（source={image_source}）")
    try:
        data = _build_pptx(plan, style_hint=p.get("style_hint") or "",
                           template_path=template_path,
                           images_by_idx=images_by_idx or None)
    except Exception as e:
        logger.exception(f"[ppt] render failed pid={pid}: {e}")
        await _touch(pid, status="outlined")
        raise HTTPException(500, f"渲染失败：{type(e).__name__}: {str(e)[:200]}")
    # 存到本地缓存目录（next 一步如果配了 storage 也可上传，但 .pptx 比较大，本地先够用）
    fname = f"ppt_{uid}_{pid}.pptx"
    fpath = os.path.join(_PPT_CACHE_DIR, fname)
    with open(fpath, "wb") as f:
        f.write(data)
    download_url = f"/api/studio/ppt/projects/{pid}/download"
    await _touch(pid, status="done", pptx_url=download_url)
    return {"ok": True, "download_url": download_url, "size_bytes": len(data)}


class PexelsKeyIn(BaseModel):
    api_key: str


@router.get("/admin/pexels-key", summary="admin 查 Pexels API key 状态")
async def admin_get_pexels_key(current_user: dict = Depends(get_current_user)):
    if (current_user.get("role") or "") != "admin":
        raise HTTPException(403, "需要 admin 权限")
    v = (await monitor_db.get_setting("pexels_api_key", "")).strip()
    return {
        "configured": bool(v),
        "preview": (v[:6] + "***" + v[-4:]) if len(v) > 12 else ("***" if v else ""),
    }


@router.put("/admin/pexels-key", summary="admin 设置 Pexels API key")
async def admin_set_pexels_key(body: PexelsKeyIn, current_user: dict = Depends(get_current_user)):
    """配置 Pexels 免费 API key（在 https://www.pexels.com/api/new/ 注册即给）。
    所有用户的 AI PPT image_source=pexels 共用这个 key。
    传空字符串 = 清除。"""
    if (current_user.get("role") or "") != "admin":
        raise HTTPException(403, "需要 admin 权限")
    key = (body.api_key or "").strip()
    await monitor_db.set_setting("pexels_api_key", key)
    return {"ok": True, "configured": bool(key)}


@router.get("/projects/{pid}/download", summary="下载渲染好的 .pptx")
async def download_project(
    pid: int, current_user: dict = Depends(get_current_user_flex),
):
    # get_current_user_flex：浏览器 <a download> 带不了 Authorization header，
    # 故额外接受 ?token=<jwt>（P0-5 修复 401 下不到文件）。
    uid = int(current_user["id"])
    p = await _get_project(pid, uid)
    if not p:
        raise HTTPException(404, "PPT 项目不存在或无权访问")
    fname = f"ppt_{uid}_{pid}.pptx"
    fpath = os.path.join(_PPT_CACHE_DIR, fname)
    if not os.path.exists(fpath):
        raise HTTPException(404, "请先调用 /render 渲染")
    safe_name = (p.get("title") or "ppt").replace("/", "_")[:60]
    return FileResponse(
        fpath,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{safe_name}.pptx",
    )
