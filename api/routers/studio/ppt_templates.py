# -*- coding: utf-8 -*-
"""AI PPT 模板复用 —— 用户上传一份 .pptx 作公司 VI 风格底版。

存储：模板 .pptx 落到磁盘 `database/ppt_templates/<id>.pptx`；行存元信息。
用法：渲染 PPT 时若 ppt_projects.template_id 非空，python-pptx 用
`Presentation(template.file_path)` 打开它，保留 master / theme / layouts，
我们只 add_slide 往 layouts[1]（一般是"标题 + 内容"版式）里塞内容。

这样用户可以一次性传"公司模板"，之后所有 PPT 都自带公司视觉风格。
"""
from __future__ import annotations

import io
import json
import logging
import os
from typing import Any, Dict, List, Optional

from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile

from ..auth import get_current_user
from ...services import monitor_db
from ...services import db as _db

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/studio/ppt/templates", tags=["AI-Studio-PPT-Templates"])
DB_PATH = monitor_db.DB_PATH

# 模板 .pptx 长期存储目录（和 monitor.db 同级）
TEMPLATES_DIR = os.path.join(os.path.dirname(DB_PATH), "ppt_templates")
os.makedirs(TEMPLATES_DIR, exist_ok=True)


def _summarize_pptx(blob: bytes) -> str:
    """从 .pptx 二进制读关键元信息生成一句话摘要给前端展示。"""
    from pptx import Presentation
    try:
        prs = Presentation(io.BytesIO(blob))
        ratio = "16:9" if abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.05 else "4:3"
        n_layouts = len(prs.slide_layouts)
        return f"{ratio} · {n_layouts} layouts · {len(prs.slides)} 张样例页"
    except Exception as exc:
        return f"(无法解析: {type(exc).__name__})"


@router.post("", summary="上传 .pptx 作风格模板")
async def upload_template(
    file: UploadFile = File(...),
    name: str = Form(""),
    current_user: dict = Depends(get_current_user),
):
    """上传一份 .pptx；存到磁盘 + 入库。不扣点。"""
    uid = int(current_user["id"])
    fname = (file.filename or "").lower()
    if not fname.endswith(".pptx"):
        raise HTTPException(400, "只支持 .pptx")
    blob = await file.read()
    if len(blob) > 30 * 1024 * 1024:
        raise HTTPException(400, "模板超过 30 MB")
    # 先解析校验
    summary = _summarize_pptx(blob)
    if summary.startswith("(无法解析"):
        raise HTTPException(400, f".pptx 解析失败：{summary}")
    pretty_name = (name.strip() or os.path.splitext(file.filename or "template")[0])[:80]

    async with _db.connect(DB_PATH) as db:
        cur = await db.execute(
            "INSERT INTO ppt_templates(user_id, name, file_path, size_bytes, layout_summary) "
            "VALUES (?,?,?,?,?)",
            (uid, pretty_name, "", len(blob), summary),
        )
        await db.commit()
        tid = cur.lastrowid
        # 现在有 id 了写文件，再回填 path
        fpath = os.path.join(TEMPLATES_DIR, f"{tid}.pptx")
        with open(fpath, "wb") as f:
            f.write(blob)
        await db.execute(
            "UPDATE ppt_templates SET file_path=? WHERE id=?", (fpath, int(tid)),
        )
        await db.commit()
    return {"ok": True, "id": tid, "name": pretty_name, "summary": summary, "size_bytes": len(blob)}


@router.get("", summary="我的 PPT 模板列表")
async def list_templates(current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    async with _db.connect(DB_PATH) as db:
        db.row_factory = _db.Row
        async with db.execute(
            "SELECT id, name, size_bytes, layout_summary, created_at "
            "FROM ppt_templates WHERE user_id=? ORDER BY id DESC", (uid,),
        ) as cur:
            rows = [dict(r) for r in await cur.fetchall()]
    return {"templates": rows}


@router.delete("/{tid}", summary="删除模板（连同磁盘文件）")
async def delete_template(tid: int, current_user: dict = Depends(get_current_user)):
    uid = int(current_user["id"])
    async with _db.connect(DB_PATH) as db:
        db.row_factory = _db.Row
        async with db.execute(
            "SELECT file_path FROM ppt_templates WHERE id=? AND user_id=?", (int(tid), uid),
        ) as cur:
            row = await cur.fetchone()
        if not row:
            raise HTTPException(404, "模板不存在")
        # 解绑还在用它的项目（不级联删，避免误删 PPT）
        await db.execute(
            "UPDATE ppt_projects SET template_id=NULL WHERE template_id=? AND user_id=?",
            (int(tid), uid),
        )
        await db.execute("DELETE FROM ppt_templates WHERE id=?", (int(tid),))
        await db.commit()
    try:
        if row["file_path"] and os.path.exists(row["file_path"]):
            os.remove(row["file_path"])
    except Exception as exc:
        logger.warning("删除模板文件失败 %s: %s", row.get("file_path"), exc)
    return {"ok": True}


# ── 内部辅助：拿模板文件路径（供 ppt.py 渲染时调用） ─────────────────────

async def get_template_path(template_id: Optional[int], user_id: int) -> Optional[str]:
    """根据 template_id + user_id 拿到磁盘路径；不存在返回 None。"""
    if not template_id:
        return None
    async with _db.connect(DB_PATH) as db:
        db.row_factory = _db.Row
        async with db.execute(
            "SELECT file_path FROM ppt_templates WHERE id=? AND user_id=?",
            (int(template_id), int(user_id)),
        ) as cur:
            row = await cur.fetchone()
    if not row:
        return None
    fp = row["file_path"]
    return fp if (fp and os.path.exists(fp)) else None
